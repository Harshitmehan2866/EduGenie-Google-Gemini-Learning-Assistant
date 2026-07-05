import os
import io
import sys

from flask import (
    Flask,
    render_template,
    request,
    jsonify,
    Response,
    stream_with_context
)

from google import genai
from dotenv import load_dotenv

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

from PIL import Image

sys.stdout.reconfigure(encoding="utf-8")

# -------------------------------------------------------
# Load Environment Variables
# -------------------------------------------------------

load_dotenv()

# -------------------------------------------------------
# Flask Configuration
# -------------------------------------------------------

app = Flask(
    __name__,
    static_folder="static",
    template_folder="templates"
)

app.secret_key = os.getenv(
    "SECRET_KEY",
    "edugenie-secret-key"
)

# -------------------------------------------------------
# Gemini Configuration
# -------------------------------------------------------

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:

    print("❌ GEMINI_API_KEY not found.")

    client = None

else:

    client = genai.Client(
        api_key=GEMINI_API_KEY
    )

# -------------------------------------------------------
# Detect Best Available Gemini Model
# -------------------------------------------------------

def get_working_model():

    if client is None:
        return "gemini-2.5-flash"

    try:

        models = list(
            client.models.list()
        )

        # Highest priority

        for model in models:

            if model.name == "models/gemini-2.5-flash":

                return model.name

        # Any 2.5 Flash model

        for model in models:

            if "gemini-2.5-flash" in model.name:

                return model.name

        # Fallback

        if models:

            return models[0].name

    except Exception as e:

        print(e)

    return "gemini-2.5-flash"

MODEL_NAME = get_working_model()

print(f"Using Model : {MODEL_NAME}")

# -------------------------------------------------------
# Extract Text From Uploaded Files
# -------------------------------------------------------

def extract_text(file):

    filename = file.filename.lower()

    try:

        if filename.endswith(".pdf"):

            reader = PdfReader(file)

            text = ""

            for page in reader.pages:

                page_text = page.extract_text()

                if page_text:

                    text += page_text + "\n"

            return text

        elif filename.endswith(".docx"):

            document = Document(file)

            return "\n".join(
                paragraph.text
                for paragraph in document.paragraphs
            )

        elif filename.endswith(".pptx"):

            presentation = Presentation(file)

            content = []

            for slide in presentation.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        content.append(shape.text)

            return "\n".join(content)

    except Exception as error:

        print(error)

    return None
# -------------------------------------------------------
# Process Uploaded Files (Image / PDF / DOCX / PPTX)
# -------------------------------------------------------

def process_multimodal(prompt, file):

    if client is None:

        return "Gemini API key is missing."

    try:

        current_model = get_working_model()

        filename = file.filename.lower()

        # ---------- Images ----------

        if filename.endswith(
            (
                ".png",
                ".jpg",
                ".jpeg",
                ".webp",
                ".heic",
                ".heif"
            )
        ):

            image = Image.open(file)

            response = client.models.generate_content(

                model=current_model,

                contents=[
                    prompt,
                    image
                ]

            )

            return response.text

        # ---------- Documents ----------

        extracted_text = extract_text(file)

        if extracted_text:

            response = client.models.generate_content(

                model=current_model,

                contents=f"""
{prompt}

Content

{extracted_text}
"""

            )

            return response.text

        return "Could not extract text."

    except Exception as error:

        return f"Gemini Error : {error}"


# -------------------------------------------------------
# Home Page
# -------------------------------------------------------

@app.route("/")
def index():

    return render_template("index.html")


# -------------------------------------------------------
# Available Gemini Models
# -------------------------------------------------------

@app.route("/models")
def models():

    if client is None:

        return jsonify({

            "error": "API key missing."

        }),500

    try:

        available=[]

        for model in client.models.list():

            available.append({

                "name":model.name

            })

        return jsonify({

            "current_model":get_working_model(),

            "models":available

        })

    except Exception as error:

        return jsonify({

            "error":str(error)

        }),500


# -------------------------------------------------------
# Ask AI
# -------------------------------------------------------

@app.route("/ask",methods=["POST"])
def ask():

    if client is None:

        return jsonify({

            "error":"API key missing."

        }),500

    data=request.get_json()

    question=data.get("question")

    language=data.get("lang","English")

    if not question:

        return jsonify({

            "error":"Question required."

        }),400

    prompt=f"""

You are EduGenie.

You are an expert educational AI assistant.

Answer the student's question clearly.

Use simple explanations.

Reply completely in {language}.

Question

{question}

"""

    try:

        response=client.models.generate_content(

            model=get_working_model(),

            contents=prompt

        )

        return jsonify({

            "answer":response.text

        })

    except Exception as error:

        return jsonify({

            "error":str(error)

        }),500


# -------------------------------------------------------
# Streaming Chat
# -------------------------------------------------------

@app.route("/chat",methods=["POST"])
def chat():

    if client is None:

        return jsonify({

            "error":"API key missing."

        }),500

    data=request.get_json()

    message=data.get("message","").strip()

    if message=="":

        return jsonify({

            "error":"Message required."

        }),400

    def generate():

        try:

            response=client.models.generate_content_stream(

                model=get_working_model(),

                contents=f"""

You are EduGenie.

Answer the following student question.

{message}

"""

            )

            for chunk in response:

                if getattr(chunk,"text",None):

                    yield f"data: {chunk.text}\n\n"

            yield "data: [DONE]\n\n"

        except Exception as error:

            yield f"data: [ERROR] {str(error)}\n\n"

    return Response(

        stream_with_context(generate()),

        mimetype="text/event-stream",

        headers={

            "Cache-Control":"no-cache",

            "Connection":"keep-alive",

            "X-Accel-Buffering":"no"

        }

    )


# -------------------------------------------------------
# Clear Chat
# -------------------------------------------------------

@app.route("/clear",methods=["POST"])
def clear_chat():

    return jsonify({

        "status":"cleared"

    })
# -------------------------------------------------------
# Quiz Generator
# -------------------------------------------------------

@app.route("/quiz", methods=["POST"])
def quiz():

    if client is None:

        return jsonify({

            "error": "API key missing."

        }), 500

    topic = request.form.get("topic")

    file = request.files.get("file")

    language = request.form.get("lang", "English")

    schema = f"""
Return ONLY valid JSON.

Do not use markdown.

Return exactly 5 questions.

Format:

[
  {{
    "question":"Question in {language}",
    "options":{{
      "A":"Option A",
      "B":"Option B",
      "C":"Option C",
      "D":"Option D"
    }},
    "correct":"A",
    "explanation":"Explanation in {language}"
  }}
]
"""

    prompt = f"""

Generate a multiple choice quiz.

{schema}

"""

    try:

        current_model = get_working_model()

        # ---------------- Topic ----------------

        if topic:

            response = client.models.generate_content(

                model=current_model,

                contents=f"{prompt}\nTopic : {topic}"

            )

            result = response.text

        # ---------------- Uploaded File ----------------

        elif file:

            result = process_multimodal(

                prompt,

                file

            )

        else:

            return jsonify({

                "error":"Topic or file required."

            }),400

        import json
        import re

        cleaned=result.strip()

        cleaned=re.sub(
            r"^```json",
            "",
            cleaned
        )

        cleaned=re.sub(
            r"^```",
            "",
            cleaned
        )

        cleaned=re.sub(
            r"```$",
            "",
            cleaned
        )

        match=re.search(

            r"\[.*\]",

            cleaned,

            re.DOTALL

        )

        if match:

            cleaned=match.group()

        quiz=json.loads(cleaned)

        return jsonify({

            "quiz":quiz

        })

    except Exception as error:

        print(error)

        return jsonify({

            "error":"Quiz generation failed.",

            "details":str(error)

        }),500


# -------------------------------------------------------
# Summarizer
# -------------------------------------------------------

@app.route("/summary", methods=["POST"])
def summary():

    if client is None:

        return jsonify({

            "error":"API key missing."

        }),500

    content=request.form.get("content")

    file=request.files.get("file")

    language=request.form.get(

        "lang",

        "English"

    )

    prompt=f"""

Provide a concise,

well structured summary

in {language}

"""

    try:

        current_model=get_working_model()

        # ------------ File ------------

        if file:

            summary=process_multimodal(

                prompt,

                file

            )

            return jsonify({

                "summary":summary

            })

        # ------------ Text ------------

        if content:

            response=client.models.generate_content(

                model=current_model,

                contents=f"""

{prompt}

{content}

"""

            )

            return jsonify({

                "summary":response.text

            })

        return jsonify({

            "error":"Text or file required."

        }),400

    except Exception as error:

        return jsonify({

            "error":str(error)

        }),500
# -------------------------------------------------------
# Health Check
# -------------------------------------------------------

@app.route("/health", methods=["GET"])
def health():

    return jsonify({

        "status": "ok",

        "service": "EduGenie",

        "model": MODEL_NAME

    })


# -------------------------------------------------------
# Error Handlers
# -------------------------------------------------------

@app.errorhandler(404)
def page_not_found(error):

    return jsonify({

        "error": "Page not found."

    }), 404


@app.errorhandler(500)
def internal_server_error(error):

    return jsonify({

        "error": "Internal Server Error."

    }), 500


# -------------------------------------------------------
# Startup Information
# -------------------------------------------------------

def startup():

    print("\n" + "=" * 60)

    print("🎓 EduGenie - Google Gemini Powered Learning Assistant")

    print("=" * 60)

    print(f"Model      : {MODEL_NAME}")

    print(f"Flask App  : Running")

    print("URL        : http://127.0.0.1:5000")

    print("=" * 60 + "\n")


# -------------------------------------------------------
# Main
# -------------------------------------------------------

if __name__ == "__main__":

    startup()

    app.run(

        host="127.0.0.1",

        port=5000,

        debug=True

    )